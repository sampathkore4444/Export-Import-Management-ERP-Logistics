#!/usr/bin/env python3
"""Builds a professional CargoFlow ERP sales presentation (.pptx) for logistics companies."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- Brand ----------
FONT = "Segoe UI"
PRIMARY = RGBColor(0x0E, 0xA5, 0xE9)
PRIMARY_DARK = RGBColor(0x02, 0x84, 0xC7)
PRIMARY_LIGHT = RGBColor(0x7D, 0xD3, 0xFC)
NAVY = RGBColor(0x0C, 0x4A, 0x6E)
INK = RGBColor(0x0F, 0x17, 0x2A)
INK_SOFT = RGBColor(0x33, 0x41, 0x55)
MUTED = RGBColor(0x64, 0x74, 0x8B)
LINE = RGBColor(0xE2, 0xE8, 0xF0)
BG_SOFT = RGBColor(0xF8, 0xFA, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BLUE = RGBColor(0xE0, 0xF2, 0xFE)
GREEN = RGBColor(0x16, 0xA3, 0x4A)
AMBER = RGBColor(0xD9, 0x77, 0x06)
DARK_BG = RGBColor(0x0B, 0x12, 0x20)

# ---------- Layout ----------
PAGE_W = 13.333
PAGE_H = 7.5
MARGIN = 0.6
CONTENT_W = PAGE_W - 2 * MARGIN  # 12.133

prs = Presentation()
prs.slide_width = Inches(PAGE_W)
prs.slide_height = Inches(PAGE_H)
BLANK = prs.slide_layouts[6]

# ---------- Helpers ----------
def add_slide():
    return prs.slides.add_slide(BLANK)


def rect(slide, x, y, w, h, color, rounded=False, radius=0.1, line=None, line_w=1.5):
    st = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    sp = slide.shapes.add_shape(st, Inches(x), Inches(y), Inches(w), Inches(h))
    if rounded:
        try:
            sp.adjustments[0] = radius
        except Exception:
            pass
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def oval(slide, x, y, w, h, color, line=None, line_w=1.5):
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = color
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp


def _set_run(r, text, size, bold, color, italic=False, font=FONT):
    r.text = text
    f = r.font
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = color
    f.name = font


def txt(slide, x, y, w, h, text, size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
        anchor=MSO_ANCHOR.TOP, font=FONT, wrap=True, line_spacing=None, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    if line_spacing:
        p.line_spacing = line_spacing
    _set_run(p.add_run(), text, size, bold, color, italic, font)
    return box


def para(tf, text, size=16, bold=False, color=INK, align=PP_ALIGN.LEFT, space_after=8,
         bullet=False, indent=0.24, font=FONT, line_spacing=None, lead=None):
    p = tf.add_paragraph()
    p.alignment = align
    if space_after is not None:
        p.space_after = Pt(space_after)
    if line_spacing:
        p.line_spacing = line_spacing
    if bullet:
        pPr = p._p.get_or_add_pPr()
        marL = int(Inches(indent))
        pPr.set("marL", str(marL))
        pPr.set("indent", str(-marL))
        buFont = pPr.makeelement(qn("a:buFont"), {"typeface": "Arial", "pitchFamily": "34", "charset": "0"})
        buChar = pPr.makeelement(qn("a:buChar"), {"char": "•"})
        pPr.append(buFont)
        pPr.append(buChar)
    if lead:
        r0 = p.add_run()
        _set_run(r0, lead, size, True, color, False, font)
    _set_run(p.add_run(), text, size, bold, color, False, font)
    return p


def footer(slide, num, total=None):
    rect(slide, MARGIN, 7.04, CONTENT_W, 0.015, LINE)
    txt(slide, MARGIN, 7.12, 6.0, 0.3, "CargoFlow ERP  ·  Confidential", 9, color=MUTED)
    label = str(num)
    if total:
        label = f"{num} / {total}"
    txt(slide, PAGE_W - MARGIN - 2.0, 7.12, 2.0, 0.3, label, 9, color=MUTED, align=PP_ALIGN.RIGHT)


def header(slide, kicker, title, num, sub=None, total=None):
    txt(slide, MARGIN, 0.40, 10.0, 0.32, kicker.upper(), 11, bold=True, color=PRIMARY_DARK)
    txt(slide, MARGIN, 0.68, CONTENT_W, 0.85, title, 28, bold=True, color=INK)
    rect(slide, MARGIN, 1.58, 1.25, 0.05, PRIMARY)
    body_top = 2.15
    if sub:
        txt(slide, MARGIN, 1.70, CONTENT_W, 0.45, sub, 14, color=MUTED)
        body_top = 2.35
    footer(slide, num, total)
    return body_top


def brand_mark(slide, x, y, size, mark_color=WHITE, bg=PRIMARY):
    sp = rect(slide, x, y, size, size, bg, rounded=True, radius=0.26)
    tf = sp.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _set_run(p.add_run(), "C", size * 30, True, mark_color)
    return sp


def bullet_card(slide, x, y, w, h, items, title=None, title_color=None, size=15,
                space_after=10, bg=None, pad=0.28, bullet=True):
    if bg is not None:
        rect(slide, x, y, w, h, bg, rounded=True, radius=0.06)
    box = slide.shapes.add_textbox(Inches(x + pad), Inches(y + pad), Inches(w - 2 * pad), Inches(h - 2 * pad))
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if title:
        para(tf, title, size=size + 3, bold=True, color=title_color or INK, space_after=8)
    for it in items:
        if isinstance(it, tuple):
            lead, rest = it
            para(tf, rest, size=size, color=INK_SOFT, space_after=space_after, bullet=bullet, lead=lead + " — ")
        else:
            para(tf, it, size=size, color=INK_SOFT, space_after=space_after, bullet=bullet)
    return box


def step_flow(slide, steps, top, gap=0.25, per_row=4, card_h=1.55, num_offset=0, color=PRIMARY, text_size=13.5):
    card_w = (CONTENT_W - (per_row - 1) * gap) / per_row
    for i, (title, desc) in enumerate(steps):
        col = i % per_row
        row = i // per_row
        x = MARGIN + col * (card_w + gap)
        y = top + row * (card_h + 0.28)
        rect(slide, x, y, card_w, card_h, WHITE, rounded=True, radius=0.08, line=LINE, line_w=1.0)
        rect(slide, x, y, card_w, 0.06, color)
        badge = rect(slide, x + 0.18, y + 0.22, 0.34, 0.34, color, rounded=True, radius=0.3)
        tf = badge.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), str(num_offset + i + 1), 13, True, WHITE)
        txt(slide, x + 0.18, y + 0.62, card_w - 0.36, 0.35, title, 14.5, bold=True, color=INK)
        txt(slide, x + 0.18, y + 0.97, card_w - 0.36, 0.5, desc, text_size, color=MUTED, line_spacing=1.05)


def highlight_panel(slide, x, y, w, h, title, lines, stat=None, stat_label=None, accent=PRIMARY):
    rect(slide, x, y, w, h, LIGHT_BLUE, rounded=True, radius=0.06)
    rect(slide, x, y, 0.09, h, accent, rounded=False)
    txt(slide, x + 0.34, y + 0.28, w - 0.6, 0.4, title, 16, bold=True, color=PRIMARY_DARK)
    box = slide.shapes.add_textbox(Inches(x + 0.34), Inches(y + 0.78), Inches(w - 0.68), Inches(h - 0.95))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    for ln in lines:
        para(tf, ln, size=13.5, color=INK_SOFT, space_after=8, bullet=True)
    if stat:
        txt(slide, x + 0.34, y + h - 1.05, w - 0.68, 0.7, stat, 30, bold=True, color=PRIMARY_DARK)
        txt(slide, x + 0.34, y + h - 0.5, w - 0.68, 0.4, stat_label, 12, color=MUTED)


# =====================================================================
# SLIDE 1 — TITLE
# =====================================================================
s = add_slide()
rect(s, 0, 0, PAGE_W, PAGE_H, NAVY)
oval(s, 9.6, -2.2, 6.5, 6.5, RGBColor(0x0E, 0x7F, 0xB0))
oval(s, 11.4, 4.6, 5.2, 5.2, RGBColor(0x0A, 0x3E, 0x5E))
oval(s, -1.6, 5.6, 4.4, 4.4, RGBColor(0x08, 0x32, 0x4D))

brand_mark(s, 0.9, 0.85, 0.85)
txt(s, 1.95, 0.95, 6.0, 0.7, "Cargo", 30, bold=True, color=WHITE)
txt(s, 3.35, 0.95, 4.0, 0.7, "Flow", 30, bold=True, color=PRIMARY_LIGHT)
txt(s, 1.95, 1.42, 8.0, 0.35, "ERP", 13, bold=True, color=RGBColor(0xBA, 0xE6, 0xFD), italic=True)

txt(s, 0.9, 2.55, 11.0, 2.1,
    "The complete operations platform for import, export & air freight",
    40, bold=True, color=WHITE, line_spacing=1.05)
txt(s, 0.9, 4.55, 10.8, 0.7,
    "Fleet  ·  Container tracking  ·  Warehouse & inventory  ·  Finance & profitability  ·  Built-in AI",
    18, color=PRIMARY_LIGHT)

rect(s, 0.9, 5.6, 3.0, 0.03, RGBColor(0x1E, 0x6F, 0x9E))
txt(s, 0.9, 5.85, 10.5, 0.45, "Prepared for:  [Logistics Company Name]", 16, bold=True, color=WHITE)
txt(s, 0.9, 6.35, 10.5, 0.35, "Built for freight forwarders, agents & 3PL operators in Cambodia and the region", 13, color=RGBColor(0x9C, 0xC8, 0xE8))

txt(s, 9.35, 6.95, 3.4, 0.35, "sales@cargoflow.app", 12, color=RGBColor(0x9C, 0xC8, 0xE8), align=PP_ALIGN.RIGHT)

# =====================================================================
# SLIDE 2 — AGENDA
# =====================================================================
s = add_slide()
body_top = header(s, "Agenda", "What we'll cover today", 2)
items = [
    ("01", "The challenge", "Where logistics operations lose time and money today"),
    ("02", "CargoFlow at a glance", "One platform from booking to closure — the big picture"),
    ("03", "Core workflows", "Import & export management with guided status pipelines"),
    ("04", "Advanced modules", "Air freight, container tracking, fleet, finance, warehouse & export docs"),
    ("05", "Built-in AI copilot", "Self-hosted AI for chat, OCR, predictions and reports"),
    ("06", "Security & getting started", "Deployment, plans, trial and next steps"),
]
gap = 0.25
col_w = (CONTENT_W - gap) / 2
for i, (num, title, desc) in enumerate(items):
    col = i % 2
    row = i // 2
    x = MARGIN + col * (col_w + gap)
    y = body_top + row * 1.55
    rect(s, x, y, col_w, 1.3, BG_SOFT, rounded=True, radius=0.08, line=LINE, line_w=1.0)
    txt(s, x + 0.28, y + 0.22, 0.9, 0.85, num, 24, bold=True, color=PRIMARY)
    txt(s, x + 1.15, y + 0.24, col_w - 1.4, 0.4, title, 16, bold=True, color=INK)
    txt(s, x + 1.15, y + 0.68, col_w - 1.4, 0.5, desc, 12.5, color=MUTED, line_spacing=1.05)

# =====================================================================
# SLIDE 3 — THE CHALLENGE
# =====================================================================
s = add_slide()
body_top = header(s, "The challenge", "Why operations teams are stretched thin", 3,
                  sub="Most forwarders run their business across a patchwork of tools — and it shows in the gaps.")
pains = [
    ("Scattered visibility", "Shipments live in spreadsheets, WhatsApp and email threads — no single source of truth."),
    ("Late alerts", "Missed ETAs, stuck jobs and expiring driver documents are discovered after the problem hits."),
    ("Weak profitability view", "Revenue, vendor costs and payments are hard to reconcile per shipment."),
    ("Siloed modes", "Sea, air and warehouse stock are tracked separately, so hand-offs get lost."),
    ("Manual documents", "Commercial invoices and packing lists are typed by hand for every shipment."),
    ("Fragile hand-offs", "Imports, exports, trucks and depots don't share one status language."),
]
gap = 0.3
col_w = (CONTENT_W - gap) / 2
card_h = 1.35
for i, (title, desc) in enumerate(pains):
    col = i % 2
    row = i // 2
    x = MARGIN + col * (col_w + gap)
    y = body_top + row * (card_h + 0.22)
    rect(s, x, y, col_w, card_h, WHITE, rounded=True, radius=0.09, line=LINE, line_w=1.0)
    rect(s, x, y, 0.08, card_h, RGBColor(0xF5, 0x96, 0x77))
    txt(s, x + 0.3, y + 0.2, col_w - 0.55, 0.4, title, 15, bold=True, color=INK)
    txt(s, x + 0.3, y + 0.62, col_w - 0.55, 0.65, desc, 12.5, color=MUTED, line_spacing=1.05)

# =====================================================================
# SLIDE 4 — AT A GLANCE (STATS)
# =====================================================================
s = add_slide()
body_top = header(s, "CargoFlow at a glance", "One platform. Every shipment. Fully tracked.", 4)
stats = [
    ("7", "Modular microservices", "Auth · Import · Fleet · Master Data · AI · Warehouse · Gateway"),
    ("80+", "API endpoints", "A complete, role-aware API behind every screen"),
    ("100%", "Self-hosted", "Runs on your own infrastructure — full data control"),
    ("3", "Languages", "English, Khmer and Chinese out of the box"),
]
card_w = (CONTENT_W - 3 * 0.25) / 4
for i, (num, title, desc) in enumerate(stats):
    x = MARGIN + i * (card_w + 0.25)
    rect(s, x, body_top, card_w, 1.75, LIGHT_BLUE, rounded=True, radius=0.08)
    txt(s, x + 0.25, body_top + 0.18, card_w - 0.5, 0.7, num, 34, bold=True, color=PRIMARY_DARK)
    txt(s, x + 0.25, body_top + 0.95, card_w - 0.5, 0.4, title, 15, bold=True, color=INK)
    txt(s, x + 0.25, body_top + 1.32, card_w - 0.5, 0.4, desc, 10.5, color=MUTED, line_spacing=1.0)

values = [
    ("Import & export", "Guided status pipelines with approvals and role-based actions"),
    ("Sea + air + warehouse", "A single source of truth across every mode"),
    ("Built-in AI", "Chat, OCR, delay prediction and weekly reports — all offline-capable"),
    ("Early-warning alerts", "Past ETA/ETD and expiring driver documents, flagged automatically"),
]
gap = 0.3
col_w = (CONTENT_W - gap) / 2
card_h = 1.15
for i, (title, desc) in enumerate(values):
    col = i % 2
    row = i // 2
    x = MARGIN + col * (col_w + gap)
    y = body_top + 2.0 + row * (card_h + 0.22)
    rect(s, x, y, col_w, card_h, WHITE, rounded=True, radius=0.09, line=LINE, line_w=1.0)
    txt(s, x + 0.3, y + 0.18, col_w - 0.55, 0.4, title, 15, bold=True, color=PRIMARY_DARK)
    txt(s, x + 0.3, y + 0.58, col_w - 0.55, 0.5, desc, 12.5, color=MUTED, line_spacing=1.05)

# =====================================================================
# SLIDE 5 — IMPORT WORKFLOW
# =====================================================================
s = add_slide()
body_top = header(s, "Core workflow  ·  01", "Import management, end to end", 5,
                  sub="Every step is a trackable status with approvals, documents and an activity log.")
import_steps = [
    ("Booking received", "Container, vessel, ETA & BL captured"),
    ("Approval & team", "Approve and assign the ops team"),
    ("License & permits", "Apply license, submit customs permit"),
    ("Truck assignment", "Internal truck or outsourced vendor"),
    ("Vessel arrival", "Record ATA at discharge port"),
    ("Customs clearance", "Inspection, duties and release"),
    ("Delivery", "Pick up, deliver, confirm unloading"),
    ("Completion", "Return container and close job"),
]
step_flow(s, import_steps, body_top, per_row=4, card_h=1.6)

# =====================================================================
# SLIDE 6 — EXPORT WORKFLOW
# =====================================================================
s = add_slide()
body_top = header(s, "Core workflow  ·  02", "Export management, end to end", 6,
                  sub="A mirrored status pipeline with its own documents and activity log.")
export_steps = [
    ("Outbound booking", "Container, vessel, ETD & BL captured"),
    ("Approval & team", "Approve and assign the ops team"),
    ("License & permits", "Export license and customs permit"),
    ("Empty pickup", "Collect empty container from depot"),
    ("Stuffing", "Confirm cargo stuffing at shipper"),
    ("Gate-in & EIR", "Record gate-in with EIR number"),
    ("Vessel departure", "Record ATD of the vessel"),
    ("Clearance & close", "Export clearance then close"),
]
step_flow(s, export_steps, body_top, per_row=4, card_h=1.6, color=PRIMARY_DARK)

# =====================================================================
# SLIDE 7 — AIR FREIGHT
# =====================================================================
s = add_slide()
body_top = header(s, "Advanced module  ·  01", "Air freight with full AWB tracking", 7,
                  sub="A dedicated air job workflow alongside your sea operations.")
air_steps = [
    ("Air booking", "AWB/HAWB, carrier, flight & route"),
    ("Approval & team", "Approve and assign ops team"),
    ("License & permits", "Air export license, customs permit"),
    ("Flight departure", "Record ATD"),
    ("Flight arrival", "Record ATA at destination"),
    ("Close", "Close job and finalize billing"),
]
step_flow(s, air_steps, body_top, per_row=4, card_h=1.6, color=RGBColor(0x8B, 0x5C, 0xF6))
txt(s, MARGIN, body_top + 3.5, CONTENT_W, 0.5,
    "Search by AWB/HAWB, carrier, origin or destination · upload documents per job · full activity history",
    13, color=MUTED)

# =====================================================================
# SLIDE 8 — CONTAINER TRACKING
# =====================================================================
s = add_slide()
body_top = header(s, "Advanced module  ·  02", "Container tracking that connects to every job", 8)
bullet_card(s, MARGIN, body_top, 7.2, 4.3, [
    ("Live registry", "Every container with unique number, size and type (dry / reefer / open-top)."),
    ("Auto event logging", "Registration and every status change are recorded automatically."),
    ("Status timeline", "Empty → loaded → in transit → arrived → returned, with event history."),
    ("In-transit view", "One click to see all containers currently on the move."),
    ("Linked to jobs", "Import and export jobs reference the same container — no double entry."),
], title="Container operations, made visible")
highlight_panel(s, 8.1, body_top, 4.63, 4.3, "Why it matters",
    ["Know exactly where every container is, without calling the port.",
     "Container numbers flow from jobs into customs documents.",
     "Reuse numbers across import and export legs."],
    stat="1 view", stat_label="of every container, everywhere")

# =====================================================================
# SLIDE 9 — FLEET MANAGEMENT
# =====================================================================
s = add_slide()
body_top = header(s, "Advanced module  ·  03", "Fleet management that keeps you compliant", 9)
bullet_card(s, MARGIN, body_top, 7.2, 4.3, [
    ("Trucks, trailers & drivers", "Full registry with availability status per asset."),
    ("Document expiry alerts", "Driver IC and licence expiries flagged 30 days out — avoid port fines."),
    ("One-click assignment", "Assign internal equipment, or outsource to an approved vendor."),
    ("Pickup scheduling", "Truck and schedule aligned to vessel ETA / ETD."),
], title="Assets you can trust, daily")
highlight_panel(s, 8.1, body_top, 4.63, 4.3, "Why it matters",
    ["Never miss a licence or ID expiry again.",
     "Internal fleet and outsourced vendors live in the same workflow.",
     "Assignments appear instantly in the job timeline."],
    stat="30 days", stat_label="advance warning on document expiry", accent=GREEN)

# =====================================================================
# SLIDE 10 — FINANCE & PROFITABILITY
# =====================================================================
s = add_slide()
body_top = header(s, "Advanced module  ·  04", "Finance & profitability per shipment", 10)
bullet_card(s, MARGIN, body_top, 7.2, 4.3, [
    ("Quotations", "Send quotes with line items; convert accepted quotes into invoices in one click."),
    ("Vendor bills & costs", "Capture vendor expenses and job costs against import or export jobs."),
    ("Payments", "Record payments and watch invoice status roll from ISSUED → PARTIAL → PAID."),
    ("Profitability", "Revenue minus costs, with margin per job."),
    ("Analytics", "30-day KPIs and your top customers by profit."),
], title="Every dollar, accounted for")
highlight_panel(s, 8.1, body_top, 4.63, 4.3, "Why it matters",
    ["Know which shipments actually make money.",
     "No more spreadsheet reconciliation between ops and accounts.",
     "Quotations flow straight into billing."],
    stat="Job-level", stat_label="profit & margin on every shipment", accent=AMBER)

# =====================================================================
# SLIDE 11 — WAREHOUSE & INVENTORY
# =====================================================================
s = add_slide()
body_top = header(s, "Advanced module  ·  05", "Warehouse & inventory, always in stock", 11)
bullet_card(s, MARGIN, body_top, 7.2, 4.3, [
    ("Multiple warehouses", "Track stock across every facility with unique warehouse codes."),
    ("Stock items", "SKU, category, unit and min-stock thresholds per item."),
    ("Auto status", "IN_STOCK / LOW_STOCK / OUT_OF_STOCK recomputed on every movement."),
    ("Stock movements", "IN and OUT with reference to the source job; insufficient stock is rejected."),
    ("Summary view", "Total items, quantities and low/out-of-stock counts at a glance."),
], title="See stock before you promise it")
highlight_panel(s, 8.1, body_top, 4.63, 4.3, "Why it matters",
    ["Know what's on hand before you quote a delivery.",
     "Low-stock alerts trigger before customers do.",
     "Movement history ties stock changes to shipments."],
    stat="0", stat_label="surprises on warehouse shelves", accent=GREEN)

# =====================================================================
# SLIDE 12 — EXPORT DOCUMENTS
# =====================================================================
s = add_slide()
body_top = header(s, "Advanced module  ·  06", "Export documents, generated in minutes", 12)
bullet_card(s, MARGIN, body_top, 7.2, 4.3, [
    ("Commercial invoices", "Line items with HS codes, auto-calculated totals per export job."),
    ("Packing lists", "Quantity, units, gross/net weight, dimensions and marks per line."),
    ("Auto-filled parties", "Shipper and consignee pulled from the linked export job."),
    ("One per job", "Each export shipment has a clean, versioned commercial invoice + packing list."),
], title="From shipment to signed paperwork")
highlight_panel(s, 8.1, body_top, 4.63, 4.3, "Why it matters",
    ["Stop retyping invoices and packing lists.",
     "HS-coded lines are audit-ready for customs.",
     "Parties stay consistent with the job record."],
    stat="1 click", stat_label="to regenerate a compliant export doc", accent=PRIMARY)

# =====================================================================
# SLIDE 13 — AI COPILOT
# =====================================================================
s = add_slide()
body_top = header(s, "Built-in AI", "Your AI copilot for logistics operations", 13,
                  sub="A private, self-hosted AI that lives with your data — no external APIs.")
ai_items = [
    ("Natural-language chat", "Ask questions about live operations in plain language — even offline."),
    ("Document OCR", "Upload a BL or PDF; the vision model extracts fields you can apply to a job."),
    ("Delay & ETA prediction", "Historical patterns produce arrival/departure forecasts with risk scores."),
    ("Weekly AI reports", "Auto-generated narrative: new jobs, revenue, bottlenecks and recommendations."),
    ("Smart job assist", "Next-step hints, stuck-job and missed-ETA anomaly detection."),
]
bullet_card(s, MARGIN, body_top, 7.2, 4.4, ai_items, title="What the copilot does")
highlight_panel(s, 8.1, body_top, 4.63, 4.4, "Built for your data",
    ["Runs 100% on your own infrastructure.",
     "All inference is local — data never leaves your network.",
     "Degrades gracefully to rule-based mode when the model is offline."],
    stat="Offline-safe", stat_label="every AI feature has a fallback", accent=RGBColor(0x8B, 0x5C, 0xF6))

# =====================================================================
# SLIDE 14 — SECURITY & DEPLOYMENT
# =====================================================================
s = add_slide()
body_top = header(s, "Trust & control", "Security, roles and deployment", 14)
bullet_card(s, MARGIN, body_top, 7.2, 4.3, [
    ("Self-hosted", "Deploy in your environment, private cloud or on-premise."),
    ("JWT security", "Token-based auth with admin, manager and staff roles."),
    ("Role-gated actions", "Approvals and destructive actions are restricted by role."),
    ("Plan-based access", "Starter, Business and Enterprise modules are enforced server-side."),
    ("Trilingual", "English, Khmer and Chinese interfaces for your whole team."),
], title="Your data, your rules")
highlight_panel(s, 8.1, body_top, 4.63, 4.3, "Compliance-ready",
    ["Data residency: everything stays inside your network.",
     "Audit-friendly activity logs on every job and document.",
     "No third-party AI vendors touching your shipments."],
    stat="100%", stat_label="data stays within your infrastructure", accent=GREEN)

# =====================================================================
# SLIDE 15 — PRICING
# =====================================================================
s = add_slide()
body_top = header(s, "Plans", "Simple plans that scale with you", 15,
                  sub="14-day free trial on every plan — no credit card required.")
plans = [
    ("Starter", "$99", "/month", ["Up to 5 users", "Import & export workflows", "Fleet management",
                                  "Master data (customers, vendors)", "Email support"], False),
    ("Business", "$249", "/month", ["Up to 25 users", "Everything in Starter",
                                    "Invoicing, quotations & vendor bills", "Documents, templates & export docs",
                                    "Finance analytics & profitability", "Containers, air freight & warehouse",
                                    "AI assistant, reports & predictions", "Priority support"], True),
    ("Enterprise", "Custom", "", ["Unlimited users", "Everything in Business",
                                  "On-premise / private-cloud hosting", "Custom integrations & workflows",
                                  "Dedicated success manager"], False),
]


def slide_text(slide, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


gap = 0.35
card_w = (CONTENT_W - 2 * gap) / 3
for i, (name, price, per, feats, featured) in enumerate(plans):
    x = MARGIN + i * (card_w + gap)
    y = body_top + 0.15
    h = 4.35
    rect(s, x, y, card_w, h, WHITE, rounded=True, radius=0.06,
         line=PRIMARY if featured else LINE, line_w=2.5 if featured else 1.0)
    if featured:
        ribbon = rect(s, x, y - 0.28, card_w, 0.5, PRIMARY, rounded=True, radius=0.3)
        tf = ribbon.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        _set_run(p.add_run(), "MOST POPULAR", 12, True, WHITE)
    txt(s, x + 0.3, y + 0.35, card_w - 0.6, 0.4, name, 18, bold=True, color=INK)
    txt(s, x + 0.3, y + 0.75, card_w - 0.6, 0.75, price, 34, bold=True, color=PRIMARY_DARK)
    if per:
        txt(s, x + 0.3, y + 1.42, card_w - 0.6, 0.35, per, 12, color=MUTED)
    box = slide_text(s, x + 0.3, y + 1.85, card_w - 0.6, h - 2.1)
    for f in feats:
        para(box, f, size=12.5, color=INK_SOFT, space_after=7, bullet=True, indent=0.2)

# =====================================================================
# SLIDE 16 — WHY CARGOFLOW (CLOSING VALUE)
# =====================================================================
s = add_slide()
body_top = header(s, "The outcome", "Why operations teams choose CargoFlow", 16)
why = [
    ("Replace the patchwork", "Spreadsheets and WhatsApp become a guided, trackable pipeline."),
    ("See problems early", "Alerts surface missed ETAs, stuck jobs and expiring documents first."),
    ("Know your real margins", "Profitability per shipment instead of end-of-year guesswork."),
    ("Scale across modes", "Sea, air and warehouse in one system as you add branches and services."),
    ("Keep data private", "Self-hosted with offline AI — no external vendor sees your shipments."),
]
bullet_card(s, MARGIN, body_top, 7.3, 4.4, why, title="From chaos to clarity")
highlight_panel(s, 8.2, body_top, 4.53, 4.4, "Measurable impact",
    ["Faster quote-to-booking turnaround",
     "Fewer missed deadlines and fines",
     "Cleaner reconciliation for finance",
     "Full audit trail on every shipment"],
    stat="End-to-end", stat_label="one system, one source of truth", accent=PRIMARY)

# =====================================================================
# SLIDE 17 — GET STARTED / CLOSING
# =====================================================================
s = add_slide()
rect(s, 0, 0, PAGE_W, PAGE_H, NAVY)
oval(s, 9.8, -2.0, 6.0, 6.0, RGBColor(0x0E, 0x7F, 0xB0))
oval(s, -1.8, 5.4, 4.6, 4.6, RGBColor(0x08, 0x32, 0x4D))
brand_mark(s, 0.9, 0.85, 0.85)
txt(s, 1.95, 0.95, 6.0, 0.7, "Cargo", 30, bold=True, color=WHITE)
txt(s, 3.35, 0.95, 4.0, 0.7, "Flow", 30, bold=True, color=PRIMARY_LIGHT)
txt(s, 1.95, 1.42, 8.0, 0.35, "ERP", 13, bold=True, color=RGBColor(0xBA, 0xE6, 0xFD), italic=True)

txt(s, 0.9, 2.45, 11.0, 1.2, "Ready to run your operations\non one platform?", 38, bold=True, color=WHITE, line_spacing=1.1)
txt(s, 0.9, 4.05, 10.5, 0.5, "Book a live demo, start a 14-day free trial, and get a guided onboarding call.", 17, color=PRIMARY_LIGHT)
rect(s, 0.9, 4.95, 3.0, 0.03, RGBColor(0x1E, 0x6F, 0x9E))

txt(s, 0.9, 5.3, 10.5, 0.45, "sales@cargoflow.app", 18, bold=True, color=WHITE)
txt(s, 0.9, 5.85, 10.5, 0.4, "Demo  ·  14-day trial  ·  Guided onboarding  ·  Trilingual (EN / KH / ZH)", 14, color=RGBColor(0x9C, 0xC8, 0xE8))
txt(s, 0.9, 6.75, 10.5, 0.35, "© 2026 CargoFlow ERP · Import, Export & Air Freight Management System", 11, color=RGBColor(0x7B, 0xA7, 0xC8))

OUT = "CargoFlow_ERP_Platform_Presentation.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
